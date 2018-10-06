import pptx
import os
import datetime
import comtypes.client as win32
import sys
from comtypes.client import Constants
import log
import logging
# with open(u'D:\\周日敬拜\\2014-10-5.pptx','rb') as file:
#     ppt_file=pptx.Presentation(file)
#     for slide in ppt_file.slides:
#         print(f'{slide.name}:{len(slide.shapes)}')
#         slide.Export(u'D:\\周日敬拜',3000,2250)
def get_logpath(path):
    """
    获取有效的存放日志路径，递归到根目录
    :param path: 路径
    :return:
    """
    if not os.path.exists(path) and os.path.split(path)[1]!='':
        get_logpath(os.path.split(path)[0])
    return os.path.split(path)[0]
def ppt_export_jpg():
    #创建日志
    logging.basicConfig(level=logging.DEBUG,
                        format='%(asctime)s %(name)-12s %(levelname)-8s %(message)s',
                        handlers=[logging.FileHandler('log_record.log','a+',encoding='utf-8')])

    FILENAME=u'D:\\周日敬拜\\2014-10-5.pptx'
    log_path=get_logpath(FILENAME)
    logging.info(FILENAME)
    #检查文件是否存在，记录报错
    if  not os.path.exists(FILENAME):
        log.log(log_path,'2014-10-5.pptx 文件不存在')
        return
    #获取本周周日日期
    today = datetime.date.today()
    add_day = 7-(today.weekday()+1)
    sunday = today + datetime.timedelta(add_day)
    #输出文件路径
    OUTPUTFILE=os.path.join(os.path.split(FILENAME)[0],f'{sunday}')
    logging.info(OUTPUTFILE)
    #创建ppt接口
    try:
        ppt=win32.CreateObject('PowerPoint.Application')
        ppt.Visible=True
        deck = ppt.Presentations.Open(FILENAME)
        # for slide in deck.Slides:
        #     print(slide.SlideNumber)
        # 导出单个ppt
        # deck.Slides[1].Export(os.path.join(OUTPUTFILE[0],'1.jpg'),'jpg',3000,2250)
        deck.Export(OUTPUTFILE,'jpg',3000,2250)
        deck.Close()
        ppt.Quit()
    except Exception as e:
        log.log(log_path,e)
        if deck is not None:
            deck.Close()
        if ppt is not None:
            ppt.Quit()
if __name__=='__main__':
    ppt_export_jpg()
